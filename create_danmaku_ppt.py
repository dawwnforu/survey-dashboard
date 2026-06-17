# -*- coding: utf-8 -*-
"""生成弹幕PPT: 弹幕墙概览 + 逐条展示 + 总结"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import textwrap, random

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
PURPLE = RGBColor(0x81, 0x8C, 0xF8)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
PINK = RGBColor(0xEC, 0x48, 0x99)

danmaku_data = [
    ("设计匹配程度系统，对不同选项设置不同权重（睡眠5权重、性格3权重），汇总得到完整匹配系数，类似雷达图", PURPLE, 18, True),
    ("建议统计打游戏类别/是否开麦，统计是否有煲电话粥的习惯，对宿舍安静环境影响很大", PURPLE, 18, True),
    ("宿舍分配急需一个更加合理并且行而有效的方案，坏的舍友可能毁了一个人的四年，好的舍友可能是未来的好兄弟好姐妹", AMBER, 22, False),
    ("可以设置一学期的适应期，如果发现室友某些方面和组宿舍之前说的不一样，可以重新组合分配宿舍", PURPLE, 18, True),
    ("改善一下隔音吧，感觉寝室里稍微动静大点周围一圈都听得见", TEXT, 16, False),
    ("喜欢的娱乐类型、未来规划、学习态度、生活习惯细则（打鼾、洗澡频率、卫生程度）都需纳入考虑", PURPLE, 18, True),
    ("建议学校宿舍分配不要唯班级/学院分配，允许自由组队。希望学校留出空房间并允许调整", TEXT, 16, False),
    ("需要安静环境，就算打游戏也不开麦，更不能用茶轴青轴等，打电话不去阳台就算了不要把自己事抖得全天下都知道", PURPLE, 18, True),
    ("尽量量化给出标准，每个人的评判标准不一样，靠感觉的话双方脑内信息不一定匹配", TEXT, 16, False),
    ("在个人作息时间基础上，考虑组宿舍自由度，最好能和相熟相知的人一起。性格爱唠嗑型 vs 保持边界感型需标注！", AMBER, 24, True),
    ("别老换宿舍，老要搬来搬去，新来的同学直接住进空出来的宿舍呗，不要在睡觉的时候装修", TEXT, 16, False),
    ("是否有过住宿经历，如若没有，是否懂得学习如何与他人共处", CYAN, 16, False),
    ("答案可以是和某些人相处好另一些人一般，不一定一个宿舍的人一概而论", TEXT, 16, False),
    ("我个人对睡觉安静比较苛刻，但室友非常体谅照顾我才导致宿舍没有矛盾，好爱他们，表白一下 ❤️", PINK, 18, False),
]

def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG
    return slide

# === Slide 1: 封面 ===
s = dark_slide()
box = s.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(12), Inches(6))
box.fill.solid(); box.fill.fore_color.rgb = CARD; box.line.fill.background()

tx = s.shapes.add_textbox(Inches(3), Inches(2.5), Inches(10), Inches(1.5))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "大学生宿舍分配研究"; p.font.size = Pt(40); p.font.bold = True
p.font.color.rgb = TEXT; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "学生真实反馈 · 弹幕墙"; p.font.size = Pt(28)
p2.font.color.rgb = PURPLE; p2.alignment = PP_ALIGN.CENTER

tx2 = s.shapes.add_textbox(Inches(3), Inches(5), Inches(10), Inches(1.5))
tf2 = tx2.text_frame
p3 = tf2.paragraphs[0]; p3.text = "来源：110份问卷调查开放题（第44-49题）| 2026年6月"
p3.font.size = Pt(16); p3.font.color.rgb = MUTED; p3.alignment = PP_ALIGN.CENTER

# === Slide 2: 弹幕墙总览 (模拟看板效果) ===
s = dark_slide()
# Semi-transparent dark overlay card
box = s.shapes.add_shape(1, Inches(0.3), Inches(0.3), Inches(15.4), Inches(8.4))
box.fill.solid(); box.fill.fore_color.rgb = CARD; box.line.fill.background()

txT = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(5), Inches(0.5))
tf = txT.text_frame
p = tf.paragraphs[0]; p.text = "📢 学生反馈弹幕墙"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = PURPLE

# Separator line
line = s.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(15), Inches(0.015))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0x33, 0x41, 0x55)
line.line.fill.background()

# Place danmaku at staggered horizontal positions (simulating scroll snapshot)
# Alternate left/right to simulate different scroll positions
random.seed(42)
tracks = [
    (0.8, 1.1), (1.2, 1.9), (0.5, 2.7), (1.5, 3.5), (0.3, 4.3),
    (1.0, 5.1), (0.7, 5.9), (0.2, 6.7), (1.3, 7.5), (0.9, 8.3),
    (0, 1.5), (1.8, 2.3), (0.6, 3.1), (0.4, 6.3),
]

for idx, (text, color, font_size, is_important) in enumerate(danmaku_data):
    y_base = tracks[idx][1] if idx < len(tracks) else idx * 0.7 + 1.0
    x_start = tracks[idx][0] if idx < len(tracks) else random.uniform(-1, 2)

    # Create text at semi-random x position to simulate scroll
    txD = s.shapes.add_textbox(Inches(x_start), Inches(y_base), Inches(14), Inches(0.6))
    tf = txD.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    # Add emoji prefix for important ones
    prefix = "💡 " if is_important else ""
    p.text = prefix + text
    p.font.size = Pt(font_size)
    p.font.bold = is_important
    p.font.color.rgb = color

# Note: user should apply Fly In animation from right in PowerPoint for dynamic effect
txN = s.shapes.add_textbox(Inches(0.5), Inches(8.3), Inches(15), Inches(0.4))
tf = txN.text_frame
p = tf.paragraphs[0]
p.text = "💡 紫色=重要建议  🟡 金色=核心诉求  ⬜ 白色=一般反馈  🩷 粉色=情感表达  🔵 青色=建议"
p.font.size = Pt(12); p.font.color.rgb = MUTED; p.alignment = PP_ALIGN.CENTER

# === Slides 3-16: 逐条展示 ===
for i, (text, color, font_size, is_important) in enumerate(danmaku_data):
    s = dark_slide()
    box = s.shapes.add_shape(1, Inches(1), Inches(1.5), Inches(14), Inches(6))
    box.fill.solid(); box.fill.fore_color.rgb = CARD; box.line.fill.background()

    # Badge
    badge = "⭐ 重要建议" if is_important else "💬 学生反馈"
    bc = AMBER if color == AMBER else (PURPLE if is_important else MUTED)
    txBadge = s.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(5), Inches(0.5))
    p = txBadge.text_frame.paragraphs[0]
    p.text = f"#{i+1} {badge}"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = bc

    # Quote text
    txQ = s.shapes.add_textbox(Inches(2), Inches(2.8), Inches(12), Inches(3.5))
    tq = txQ.text_frame; tq.word_wrap = True
    p = tq.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size + 8); p.font.bold = (color == AMBER)
    p.font.color.rgb = color; p.line_spacing = Pt(44)

    # Source
    txS = s.shapes.add_textbox(Inches(2), Inches(6.8), Inches(12), Inches(0.5))
    p = txS.text_frame.paragraphs[0]
    p.text = "—— 问卷开放题 #44-49"; p.font.size = Pt(14)
    p.font.color.rgb = MUTED; p.font.italic = True; p.alignment = PP_ALIGN.RIGHT

# === Last Slide: 总结 ===
s = dark_slide()
box = s.shapes.add_shape(1, Inches(2), Inches(1), Inches(12), Inches(7))
box.fill.solid(); box.fill.fore_color.rgb = CARD; box.line.fill.background()

txT = s.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(11), Inches(0.8))
p = txT.text_frame.paragraphs[0]
p.text = "📊 学生反馈核心诉求 TOP 5"; p.font.size = Pt(32); p.font.bold = True
p.font.color.rgb = TEXT; p.alignment = PP_ALIGN.CENTER

points = [
    ("① 作息匹配是第一刚需", "睡眠安静权重最高，被反复提及"),
    ("② 游戏/电话噪音 → 细化分类", "是否开麦、键盘类型、打电话习惯"),
    ("③ 量化标准 + 个人简介系统", "主观感觉不可靠，需要客观匹配系数"),
    ("④ 设置适应期 + 允许重分配", "说的和实际不一样，需要容错机制"),
    ("⑤ 性格维度标注", "社牛型/边界感型，避免社交偏好冲突"),
]

for i, (title, desc) in enumerate(points):
    y_pos = 2.8 + i * 1.0
    # Title
    tx1 = s.shapes.add_textbox(Inches(3), Inches(y_pos), Inches(10), Inches(0.5))
    p = tx1.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = [AMBER, PURPLE, GREEN, CYAN, PINK][i]

    # Description
    tx2 = s.shapes.add_textbox(Inches(3.3), Inches(y_pos + 0.4), Inches(9), Inches(0.4))
    p = tx2.text_frame.paragraphs[0]
    p.text = desc; p.font.size = Pt(14); p.font.color.rgb = MUTED

out_path = r'D:\HuaweiMoveData\Users\xiaoy\Desktop\社会调查实践\学生反馈弹幕墙.pptx'
prs.save(out_path)
print(f"Done! {len(prs.slides)} slides saved.")

print("""
使用方法:
  Slide 1: 封面
  Slide 2: 弹幕墙总览（14条弹幕散布）
  Slide 3-16: 逐条展示（共14条）
  Slide 17: 核心诉求总结

  若要动态滚动效果:
  PowerPoint → 选中Slide 2的文本框 → 动画 → 飞入（自右侧）
  → 对每条设置不同延迟（0s, 0.5s, 1s...）即可实现弹幕滚动效果
""")
